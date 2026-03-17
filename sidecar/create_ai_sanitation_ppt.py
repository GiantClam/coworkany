#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI在环卫中的应用及2026年展望 - PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

# 1. 封面
add_title_slide(prs, 
    "AI在环卫中的应用及2026年展望",
    "智慧环卫：技术驱动的城市清洁革命"
)

# 2. 环卫行业定义和发展背景
add_content_slide(prs, 
    "1. 环卫行业定义与发展背景",
    [
        "• 环卫行业定义",
        "  - 城市环境卫生管理的核心产业",
        "  - 涵盖道路清扫、垃圾收运、分类处理等全链条服务",
        "  - 2025年中国城市环卫市场规模达2315亿元",
        "",
        "• 发展背景",
        "  - 双碳目标推动绿色低碳转型",
        "  - 智慧城市建设加速数字化升级",
        "  - 人工成本上升倒逼智能化改造",
        "  - 国务院《关于深入实施'人工智能+'行动的意见》政策支持"
    ]
)

# 3. AI在环卫中的主要应用场景
add_content_slide(prs,
    "2. AI在环卫中的主要应用场景",
    [
        "• 智能环卫车辆",
        "  - L4级无人驾驶清扫车：24小时自主作业，减少人力需求",
        "  - AI感知与导航系统：多模态传感器融合、高精度定位",
        "  - 零排放电动动力：绿色环保，降低运营成本",
        "",
        "• 智能垃圾分拣",
        "  - AI视觉识别技术：识别20+种垃圾类型（塑料、金属、纸张等）",
        "  - 机器人分拣：每分钟70件，速度是人工2倍，可24小时工作",
        "  - 案例：联运环境AI分拣系统、北京西红门智能分拣机器人",
        "",
        "• 智慧调度系统",
        "  - AI预测垃圾产生量：分析历史数据、天气、节假日等因素",
        "  - 优化收运路线：杭州某区通过AI调度减少20%环卫车空驶率",
        "  - 无人机智慧巡检：大理市实现环卫行业首个无人机巡检应用"
    ]
)

# 4. 2026年最新技术和趋势
add_content_slide(prs,
    "3. 2026年最新技术与趋势",
    [
        "• 技术突破",
        "  - 6G+AI全场景应用：南京紫金山科技城落地全球首个6G+AI无人环卫",
        "  - 纯视觉自动驾驶：不依赖雷达，降低成本提升普及率",
        "  - 物联网深度融合：路径规划、故障报警、远程监控一体化",
        "",
        "• 市场趋势",
        "  - 市场规模：全球AI清洁机器人市场2026年将超85亿美元（CAGR 17.9%）",
        "  - 项目爆发：2025年国内成功开标无人环卫项目超220项",
        "  - 无人环卫市场价值预计超3000亿元",
        "",
        "• 政策驱动",
        "  - L3以上自动驾驶汽车加速商用化",
        "  - 智能环卫机器人成为智慧城市核心装备",
        "  - 数字化转型成为环卫企业必选项"
    ]
)

# 5. 成功案例分析
add_content_slide(prs,
    "4. 成功案例分析",
    [
        "• 案例一：九爪智能环卫综合体（智慧城市大会标杆项目）",
        "  - AI分选系统实现垃圾'变废为宝'",
        "  - 前瞻性设计理念与硬核AI技术结合",
        "",
        "• 案例二：杭州AI调度系统",
        "  - 通过AI预测和优化，减少20%环卫车空驶率",
        "  - 显著降低运营成本，提升服务效率",
        "",
        "• 案例三：驭势科技L4级无人环卫车",
        "  - 整合车规级域控制器与领先AI算法",
        "  - 重塑劳动力价值，开启智能清扫新纪元",
        "",
        "• 案例四：景德镇玉禾田'四精'管理模式",
        "  - 入选2024年度数字化转型驱动环卫管理创新案例",
        "  - 数字化运营实现精细化管理"
    ]
)

# 6. 未来展望
add_content_slide(prs,
    "5. 未来展望",
    [
        "• 技术演进方向",
        "  - 从'人工主导'到'人机协同'再到'无人智慧化'",
        "  - AI大模型赋能：更强的环境理解和决策能力",
        "  - 多技术融合：AI+IoT+5G/6G+边缘计算+区块链",
        "",
        "• 产业发展趋势",
        "  - 环卫机器人从'试点'走向'规模化商用'",
        "  - 传统环卫企业加速数字化转型",
        "  - 新能源+智能化成为行业标配",
        "",
        "• 社会价值",
        "  - 提升城市环境质量，助力智慧城市建设",
        "  - 降低环卫工人劳动强度，改善工作环境",
        "  - 推动绿色低碳发展，实现可持续运营",
        "",
        "• 挑战与机遇",
        "  - 技术成熟度与成本平衡",
        "  - 政策法规配套完善",
        "  - 产业链协同创新"
    ]
)

# 7. 结束页
add_title_slide(prs,
    "谢谢观看",
    "AI赋能环卫，共创智慧未来"
)

# 保存文件
prs.save('AI在环卫中的应用及2026年展望.pptx')
print("✅ PPT创建成功：AI在环卫中的应用及2026年展望.pptx")

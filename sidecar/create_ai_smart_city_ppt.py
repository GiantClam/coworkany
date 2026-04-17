#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成"2026年AI在智慧城市中的发展"PPT演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.templates[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # 设置标题格式
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_content_slide(prs, title, content_items):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.templates[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(12)

def main():
    # 创建演示文稿对象
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 标题页
    create_title_slide(
        prs,
        "2026年AI在智慧城市中的发展",
        "人工智能驱动的城市未来\n2026年最新趋势与应用"
    )
    
    # 2. 智慧城市的定义和发展背景
    create_content_slide(
        prs,
        "智慧城市：定义与发展背景",
        [
            "定义：利用物联网、大数据、人工智能等技术，实现城市管理和服务的智能化",
            "发展背景：",
            "  • 全球城市化进程加速，城市人口持续增长",
            "  • 城市管理面临交通拥堵、环境污染、资源短缺等挑战",
            "  • AI技术成熟度提升，从愿景走向实际应用",
            "  • 数字化转型成为城市发展的核心战略",
            "市场规模：AI智慧城市应用市场近年来快速增长，2026年持续扩张"
        ]
    )
    
    # 3. AI在智慧城市中的主要应用场景
    create_content_slide(
        prs,
        "AI在智慧城市中的主要应用场景",
        [
            "1. 智能交通管理",
            "  • 实时交通流量监控与优化",
            "  • AI驱动的信号灯控制系统",
            "  • 智能停车引导与管理",
            "",
            "2. 智慧能源管理",
            "  • 电网负载预测与优化",
            "  • 长时储能系统（应对AI数据中心需求）",
            "  • 可再生能源智能调度"
        ]
    )
    
    # 4. AI应用场景（续）
    create_content_slide(
        prs,
        "AI应用场景（续）",
        [
            "3. 城市基础设施监控",
            "  • 实时传感器网络监测",
            "  • 预测性维护系统",
            "  • 数字孪生技术模拟城市运行",
            "",
            "4. 认知政府服务",
            "  • 数字化采购与工作流自动化",
            "  • GIS地理信息系统集成",
            "  • AI辅助决策支持系统"
        ]
    )
    
    # 5. 2026年最新技术和趋势
    create_content_slide(
        prs,
        "2026年最新技术和趋势",
        [
            "1. 认知能力加速发展",
            "  • 高级传感技术：实时城市感知",
            "  • 数字孪生：虚拟城市模拟与优化",
            "  • 智能体系统：自主决策与执行",
            "",
            "2. 运营成熟度提升",
            "  • 从概念验证到规模化部署",
            "  • 跨部门数据集成与共享",
            "  • 标准化与互操作性增强"
        ]
    )
    
    # 6. 2026年趋势（续）
    create_content_slide(
        prs,
        "2026年技术趋势（续）",
        [
            "3. AI基础设施扩张",
            "  • 数据中心容量大幅增长（如德国计划2030年翻倍）",
            "  • AI算力需求推动能源存储创新",
            "  • 边缘计算与云计算协同",
            "",
            "4. 可持续发展导向",
            "  • 绿色AI技术应用",
            "  • 碳排放监测与管理",
            "  • 循环经济智能化支持"
        ]
    )
    
    # 7. 成功案例分析
    create_content_slide(
        prs,
        "成功案例分析",
        [
            "案例1：中国雄安新区",
            "  • 定位：未来之城示范项目",
            "  • 特点：从规划阶段就融入AI和智慧城市理念",
            "  • 应用：智能基础设施、数字孪生城市管理",
            "",
            "案例2：全球智慧城市奖项（2026）",
            "  • 趋势：项目从愿景转向实际影响",
            "  • 成果：数字采购、实时监控、数据集成广泛应用",
            "  • 特点：运营成熟度显著提升"
        ]
    )
    
    # 8. 成功案例（续）
    create_content_slide(
        prs,
        "成功案例分析（续）",
        [
            "案例3：智能交通系统",
            "  • 实时交通重新布线技术",
            "  • AI优化信号灯系统，减少拥堵30-40%",
            "  • 预测性交通管理，提前应对高峰",
            "",
            "案例4：能源管理创新",
            "  • 长时储能系统部署（应对AI数据中心需求）",
            "  • 智能电网负载平衡",
            "  • 可再生能源利用率提升"
        ]
    )
    
    # 9. 挑战与机遇
    create_content_slide(
        prs,
        "面临的挑战与机遇",
        [
            "挑战：",
            "  • 数据隐私与安全保护",
            "  • 跨部门协调与数据共享障碍",
            "  • 技术投资回报周期长",
            "  • AI伦理与治理问题",
            "",
            "机遇：",
            "  • 提升城市运营效率20-30%",
            "  • 改善居民生活质量",
            "  • 创造新的经济增长点",
            "  • 推动可持续发展目标实现"
        ]
    )
    
    # 10. 未来展望
    create_content_slide(
        prs,
        "未来展望：2026年及以后",
        [
            "短期展望（2026-2027）：",
            "  • AI应用从试点走向规模化部署",
            "  • 数字孪生技术成为城市管理标配",
            "  • 跨城市数据共享与协作增强",
            "",
            "中长期展望（2028-2030）：",
            "  • 全自主智能城市系统出现",
            "  • 人机协同治理模式成熟",
            "  • 智慧城市成为碳中和关键路径",
            "  • AI驱动的城市创新生态系统形成"
        ]
    )
    
    # 11. 结论页
    create_content_slide(
        prs,
        "结论",
        [
            "AI正在重塑城市的运作方式",
            "",
            "2026年标志着智慧城市从愿景到实际影响的转折点",
            "",
            "技术成熟度、运营经验和生态系统的完善，为未来发展奠定基础",
            "",
            "智慧城市不仅是技术创新，更是治理模式和生活方式的变革",
            "",
            "未来属于那些能够有效整合AI技术、以人为本的城市"
        ]
    )
    
    # 12. 致谢页
    create_content_slide(
        prs,
        "谢谢！",
        [
            "感谢您的聆听",
            "",
            "Questions & Discussion",
            "",
            "基于2026年最新研究与行业报告",
            "数据来源：IDC、Deloitte、行业分析报告"
        ]
    )
    
    # 保存文件
    prs.save('AI智慧城市发展_2026.pptx')
    print("PPT文件已成功创建：AI智慧城市发展_2026.pptx")

if __name__ == "__main__":
    main()

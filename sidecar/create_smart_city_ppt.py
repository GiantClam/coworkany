from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

prs = Presentation()

# Theme colors
TITLE_COLOR = RGBColor(20, 45, 100)
ACCENT_COLOR = RGBColor(0, 112, 192)
TEXT_COLOR = RGBColor(50, 50, 50)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    t = slide.shapes.title.text_frame.paragraphs[0].runs[0].font
    t.size = Pt(38)
    t.bold = True
    t.color.rgb = TITLE_COLOR


def add_bullets_slide(title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0].font
    title_run.color.rgb = TITLE_COLOR
    title_run.bold = True

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        if isinstance(b, tuple):
            p.text = b[0]
            p.level = b[1]
        else:
            p.text = b
            p.level = 0
        p.font.size = Pt(22 if p.level == 0 else 18)
        p.font.color.rgb = TEXT_COLOR

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_two_content_slide(title, left_title, left_points, right_title, right_points, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = TITLE_COLOR

    left = slide.shapes.placeholders[1].text_frame
    left.clear()
    left.paragraphs[0].text = left_title
    left.paragraphs[0].font.bold = True
    left.paragraphs[0].font.size = Pt(22)
    left.paragraphs[0].font.color.rgb = ACCENT_COLOR
    for pt in left_points:
        p = left.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR

    right = slide.shapes.placeholders[2].text_frame
    right.clear()
    right.paragraphs[0].text = right_title
    right.paragraphs[0].font.bold = True
    right.paragraphs[0].font.size = Pt(22)
    right.paragraphs[0].font.color.rgb = ACCENT_COLOR
    for pt in right_points:
        p = right.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


now = datetime.now().strftime('%Y-%m-%d')

# 1 Cover
add_title_slide(
    "2026年AI在智慧城市中的发展",
    f"趋势研判、应用场景与案例分析\n汇报日期：{now}"
)

# 2 Agenda
add_bullets_slide(
    "目录",
    [
        "1. 智慧城市的定义和发展背景",
        "2. AI在智慧城市中的主要应用场景",
        "3. 2026年最新技术和趋势",
        "4. 成功案例分析",
        "5. 未来展望与行动建议",
    ]
)

# 3 Definition & background
add_bullets_slide(
    "1. 智慧城市的定义与发展背景",
    [
        "智慧城市：以数据、连接、算法和治理协同提升城市运行效率与居民福祉。",
        "从“数字化”走向“智能化”：IoT感知层 + 云边协同 + AI决策层成为主架构。",
        "评估维度从技术能力扩展到“以人为本”：生活质量、包容性、信任与透明度。",
        "IMD Smart City Index持续强调“技术+人文”双平衡，治理与公众信任成为核心竞争力。",
    ],
    notes="证据来源：IMD Smart City Index 2025/2026 页面（https://www.imd.org/smart-city-observatory/Home/ ; https://www.imd.org/smart-city-observatory/smart-city-index）"
)

# 4 Scenarios
add_two_content_slide(
    "2. AI在智慧城市中的主要应用场景",
    "城市运行与基础设施",
    [
        "智能交通信号优化与拥堵预测",
        "能源负荷预测与电网调度",
        "供水/管网异常检测与预测性维护",
        "城市安防与应急联动决策",
    ],
    "公共服务与民生",
    [
        "政务智能客服与多语种服务",
        "医疗资源调度与分级诊疗辅助",
        "教育与社区服务的精准供给",
        "空气质量、噪声与热岛治理",
    ],
    notes="参考：Smart Cities World 2026（AI与城市运营/出行）；LTA新加坡交通管理实践。"
)

# 5 2026 trends
add_bullets_slide(
    "3. 2026年最新技术与趋势",
    [
        "趋势1：生成式AI与城市数字孪生融合，支持“仿真—推演—决策”闭环。",
        "趋势2：边缘AI加速落地，交通、安防、能源场景强调低时延与本地推理。",
        "趋势3：从“屏幕AI”走向“物理AI”，机器人与自动系统进入城市空间。",
        "趋势4：AI治理成为新型基础设施：模型审计、实时监测、责任追溯。",
        "趋势5：跨部门数据协同与AI代理（Agentic AI）推动流程自动化。",
    ],
    notes="证据：\n1) Smart Cities World, 2026-02-19: https://www.smartcitiesworld.net/ai-and-machine-learning/how-ai-will-define-cities-and-mobility-in-2026\n2) arXiv综述（GenAI+Urban Digital Twins）: https://arxiv.org/abs/2405.19464v2\n3) WEF 2026（human-centred physical AI / governance）: https://www.weforum.org/stories/2026/02/human-centred-physical-ai-transforming-cities/ ; https://www.weforum.org/stories/2026/02/why-governance-is-the-new-infrastructure-for-physical-ai/"
)

# 6 Case 1 Singapore
add_bullets_slide(
    "4. 成功案例分析（1）新加坡：AI交通优化",
    [
        "问题：高密度城市交通需要兼顾效率与安全。",
        "做法：基于GLIDE系统与实时交通数据，动态优化信号配时。",
        "价值：减少无效等待时间，提升路网通行效率并保持安全优先。",
        "启示：先建立“数据基础设施+规则引擎”，再迭代AI增强。",
    ],
    notes="来源：LTA Media Reply（2026-03-12）https://www.lta.gov.sg/content/ltagov/en/newsroom/2026/3/media-replies/traffic-data-used-to-optimise-network-flow-while-prioritising-sa.html"
)

# 7 Case 2 Barcelona
add_bullets_slide(
    "4. 成功案例分析（2）巴塞罗那：城市数字孪生与空气质量",
    [
        "问题：空气质量与交通排放管理复杂，需跨系统联动。",
        "做法：构建城市空气质量数据库，推进数字孪生能力建设。",
        "价值：支持政策情景模拟，提升环境治理的可解释性与前瞻性。",
        "启示：数字孪生应从“高价值子系统”切入，逐步扩展全城模型。",
    ],
    notes="来源：UrbanAIR项目更新（2026）https://www.urbanair-project.eu/post/bsc-develops-a-database-to-advance-towardsa-digital-twin-of-air-quality-in-barcelona"
)

# 8 Case 3 Governance
add_bullets_slide(
    "4. 成功案例分析（3）治理实践：以信任为核心的AI城市治理",
    [
        "问题：AI系统规模化后，算法偏差、问责与隐私风险同步上升。",
        "做法：引入治理框架（透明度、审计、持续监测、公众沟通）。",
        "价值：降低政策与技术摩擦，提升市民采纳度与跨部门协同效率。",
        "启示：‘治理能力’正在成为智慧城市竞争力的重要组成。",
    ],
    notes="来源：WEF 2026治理相关文章 + IMD对透明与信任的强调。"
)

# 9 Future outlook
add_bullets_slide(
    "5. 未来展望（2026-2030）",
    [
        "城市AI将从“单点智能”走向“城市级智能体协同”。",
        "‘数字孪生+实时数据+多智能体’将成为城市运营中枢。",
        "公共价值导向将超越单纯效率：公平、韧性、可持续成为KPI。",
        "法规与标准将加速完善，模型合规与数据主权要求更高。",
        "建议：构建“场景优先、治理先行、渐进式扩展”的实施路线图。",
    ]
)

# 10 Closing
add_bullets_slide(
    "结论与行动建议",
    [
        "短期（0-12个月）：聚焦交通/能源/政务三类高ROI场景试点。",
        "中期（1-2年）：打通跨部门数据资产，建设统一AI治理与评估体系。",
        "长期（3年以上）：构建城市级数字孪生与智能体运营平台。",
        "关键成功因素：数据质量、组织协同、治理透明、持续迭代。",
    ]
)

output = "2026年AI在智慧城市中的发展_专业版.pptx"
prs.save(output)
print(output)

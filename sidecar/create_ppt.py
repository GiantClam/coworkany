#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建AI智慧城市PPT演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.templates[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # 设置标题格式
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_content_slide(prs, title, content_list):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.templates[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(6)

def create_two_column_slide(prs, title, left_content, right_content):
    """创建两栏内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    tf_left = left_box.text_frame
    for item in left_content:
        p = tf_left.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # 右栏
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.5), Inches(5))
    tf_right = right_box.text_frame
    for item in right_content:
        p = tf_right.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(6)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 第1页：标题页
    create_title_slide(prs, 
                      "2026年AI在智慧城市中的发展",
                      "技术创新 · 场景应用 · 未来展望")
    
    # 第2页：目录
    create_content_slide(prs, "目录", [
        "1. 智慧城市的定义和发展背景",
        "2. AI在智慧城市中的主要应用场景",
        "3. 2026年最新技术和趋势",
        "4. 成功案例分析",
        "5. 未来展望"
    ])
    
    # 第3页：智慧城市定义
    create_content_slide(prs, "智慧城市的定义", [
        "核心概念：利用物联网、云计算、大数据、人工智能等新一代信息技术",
        "目标：提升城市治理效率、改善民生服务、促进可持续发展",
        "关键要素：",
        "  • 智能基础设施",
        "  • 数据驱动决策",
        "  • 市民参与互动",
        "  • 跨部门协同"
    ])
    
    # 第4页：发展背景
    create_content_slide(prs, "发展背景", [
        "全球城镇化趋势：2026年预计全球60%以上人口居住在城市",
        "技术成熟度：AI、5G/6G、边缘计算技术日趋成熟",
        "政策支持：各国政府大力推动数字化转型",
        "市场规模：全球智慧城市市场规模预计突破3万亿美元"
    ])
    
    # 第5页：AI应用场景概览
    create_content_slide(prs, "AI在智慧城市中的主要应用场景", [
        "1. 智能交通管理",
        "2. 公共安全与应急管理",
        "3. 智慧能源与环境",
        "4. 智慧医疗与健康",
        "5. 智慧政务与服务"
    ])
    
    # 第6页：智能交通管理
    create_content_slide(prs, "应用场景1：智能交通管理", [
        "实时交通优化：AI分析交通流量，动态调整信号灯",
        "自动驾驶：无人驾驶公交、出租车服务",
        "停车管理：智能停车引导和预约系统",
        "预测性维护：道路和基础设施状态监测"
    ])
    
    # 第7页：公共安全
    create_content_slide(prs, "应用场景2：公共安全与应急管理", [
        "智能监控：异常行为识别、人群密度监测",
        "应急响应：AI辅助灾害预警和资源调度",
        "犯罪预防：数据分析预测高风险区域",
        "消防安全：火灾风险评估和快速响应"
    ])
    
    # 第8页：智慧能源与环境
    create_content_slide(prs, "应用场景3：智慧能源与环境", [
        "智能电网：需求预测和负载平衡",
        "环境监测：空气质量、噪音、水质实时监控",
        "节能优化：建筑能耗智能管理",
        "碳排放管理：碳足迹追踪和减排建议"
    ])
    
    # 第9页：2026年核心技术
    create_content_slide(prs, "2026年核心技术突破", [
        "1. 多模态大模型应用",
        "   • 城市级GPT：理解和处理城市多源数据",
        "   • 视觉-语言融合：监控视频智能分析",
        "2. 边缘AI与分布式计算",
        "   • 实时处理：毫秒级响应的边缘智能",
        "   • 隐私保护：数据本地化处理",
        "3. 数字孪生城市",
        "   • 虚实映射：城市全要素数字化建模",
        "4. 联邦学习与隐私计算",
        "   • 数据安全：不共享原始数据的协同学习"
    ])
    
    # 第10页：发展趋势
    create_content_slide(prs, "2026年发展趋势", [
        "从单点应用到系统集成",
        "  打通各垂直领域，实现全域智能化",
        "从被动响应到主动预测",
        "  预测性分析成为标配",
        "从技术驱动到场景驱动",
        "  更关注实际问题解决和用户体验",
        "从中心化到分布式",
        "  边缘智能与云端协同"
    ])
    
    # 第11页：案例1 - 新加坡
    create_content_slide(prs, "成功案例1：新加坡智慧国家计划", [
        "背景：全球智慧城市领先者",
        "AI应用：",
        "  • 智能交通系统：减少30%拥堵时间",
        "  • 智慧组屋：能耗降低20%",
        "  • 虚拟助手：24/7政务服务",
        "成效：",
        "  • 政府效率提升40%",
        "  • 市民满意度达85%"
    ])
    
    # 第12页：案例2 - 杭州
    create_content_slide(prs, "成功案例2：中国杭州城市大脑", [
        "背景：阿里云支持的城市级AI平台",
        "AI应用：",
        "  • 交通信号优化：通行时间减少15.3%",
        "  • 120急救调度：响应时间缩短50%",
        "  • 城市事件处理：自动识别和派单",
        "成效：",
        "  • 覆盖420平方公里",
        "  • 日均处理事件8万余起"
    ])
    
    # 第13页：案例3 - 迪拜
    create_content_slide(prs, "成功案例3：迪拜智慧城市2026", [
        "背景：中东智慧城市先锋",
        "AI应用：",
        "  • 无人驾驶出租车：占比25%",
        "  • 区块链+AI政务：100%无纸化",
        "  • 智能建筑：全城绿色建筑认证",
        "成效：",
        "  • 政府交易100%数字化",
        "  • 旅行时间节省2500万小时/年"
    ])
    
    # 第14页：未来展望 - 技术演进
    create_content_slide(prs, "未来展望：技术演进方向", [
        "短期（2026-2028）",
        "  • AI普及化：中小城市广泛应用",
        "  • 标准统一：技术标准逐步统一",
        "中期（2028-2030）",
        "  • AGI应用：通用人工智能试点",
        "  • 全域感知：城市全要素实时感知",
        "长期（2030+）",
        "  • 城市智能体：自我学习、自我优化",
        "  • 人机共生：AI深度融入日常生活"
    ])
    
    # 第15页：挑战与建议
    create_two_column_slide(prs, "面临的挑战与发展建议",
        [
            "挑战：",
            "• 数据安全与隐私保护",
            "• 技术伦理与公平性",
            "• 投资回报与商业模式",
            "• 跨部门协同与数据共享",
            "• 人才培养与技术更新"
        ],
        [
            "建议：",
            "对政府：",
            "• 制定长期规划和配套政策",
            "• 加强数据治理和标准建设",
            "对企业：",
            "• 聚焦场景化解决方案",
            "• 加强技术研发和生态合作",
            "对市民：",
            "• 提升数字素养",
            "• 积极参与智慧城市建设"
        ]
    )
    
    # 第16页：结语
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = "结语"
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "2026年，AI技术在智慧城市中的应用已从概念走向深度实践。"
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "通过技术创新、政策支持和多方协同，智慧城市正在让城市生活更安全、更高效、更宜居。"
    p2.font.size = Pt(20)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "未来已来，智慧城市的黄金时代正在开启！"
    p3.font.size = Pt(24)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(0, 102, 204)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)
    
    # 保存PPT
    prs.save('2026年AI在智慧城市中的发展.pptx')
    print("PPT创建成功！文件名：2026年AI在智慧城市中的发展.pptx")

if __name__ == "__main__":
    main()
